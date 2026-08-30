import os
import re
import json
import logging
from typing import List, Dict, Any, Optional
from dataclasses import dataclass

logger = logging.getLogger("aetheris.rag")

@dataclass
class DocumentChunk:
    chunk_id: str
    doc_id: str
    chapter: str
    content: str
    page_number: int


class RAGEngine:
    """
    Multimodal RAG & Knowledge Grounding Engine.
    Parses educational materials (PDF, DOCX, PPTX, TXT), indexes them into
    semantic chunks, and supports cross-lingual retrieval and topic knowledge synthesis.
    """
    def __init__(self, persist_dir: str = "./chroma_db"):
        self.persist_dir = persist_dir
        self.in_memory_docs: Dict[str, Dict[str, Any]] = {}
        self.in_memory_chunks: Dict[str, List[DocumentChunk]] = {}
        os.makedirs(persist_dir, exist_ok=True)
        logger.info("RAGEngine initialized with persistent and memory stores.")

    def parse_document(self, file_path: str, doc_id: str) -> Dict[str, Any]:
        """
        Parses PDF, DOCX, PPTX, or TXT file and extracts structured text, chapters, and sections.
        """
        ext = os.path.splitext(file_path)[1].lower()
        extracted_text = ""
        chapters = []
        pages_count = 1

        try:
            if ext == ".pdf":
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    pages_count = len(reader.pages)
                    for idx, page in enumerate(reader.pages):
                        page_text = page.extract_text() or ""
                        extracted_text += f"\n--- Page {idx+1} ---\n" + page_text
                except Exception as pdf_err:
                    logger.warning(f"pypdf reader fallback: {pdf_err}")
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        extracted_text = f.read()

            elif ext == ".docx":
                try:
                    import docx
                    doc = docx.Document(file_path)
                    extracted_text = "\n".join([p.text for p in doc.paragraphs if p.text])
                    pages_count = max(1, len(extracted_text) // 1800)
                except Exception as docx_err:
                    logger.warning(f"docx reader fallback: {docx_err}")
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        extracted_text = f.read()

            elif ext == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    pages_count = len(prs.slides)
                    for idx, slide in enumerate(prs.slides):
                        slide_texts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_texts.append(shape.text)
                        extracted_text += f"\n--- Slide {idx+1} ---\n" + "\n".join(slide_texts)
                except Exception as pptx_err:
                    logger.warning(f"pptx reader fallback: {pptx_err}")
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        extracted_text = f.read()

            else:  # .txt, .md, etc.
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    extracted_text = f.read()
                pages_count = max(1, len(extracted_text) // 1500)

            # Heuristic extraction of chapters / headings
            heading_matches = re.findall(r"(?:Chapter|Section|\b[0-9]+\.[0-9]+|\bUnit)\s*[:\-]?\s*([^\n\r]+)", extracted_text, re.IGNORECASE)
            if heading_matches:
                chapters = list(dict.fromkeys([h.strip() for h in heading_matches[:8]]))
            else:
                chapters = ["Overview & Core Definitions", "Fundamental Principles", "Practical Examples & Applications", "Summary & Key Formulas"]

            # Chunk document
            chunks = self._chunk_text(extracted_text, doc_id, chapters)
            self.in_memory_docs[doc_id] = {
                "file_path": file_path,
                "chapters": chapters,
                "total_pages": pages_count,
                "full_text": extracted_text
            }
            self.in_memory_chunks[doc_id] = chunks

            logger.info(f"Parsed doc '{doc_id}' into {len(chunks)} semantic chunks across {pages_count} pages.")
            return {
                "doc_id": doc_id,
                "total_pages": pages_count,
                "chapters": chapters,
                "chunk_count": len(chunks)
            }

        except Exception as e:
            logger.error(f"Error parsing document {file_path}: {e}")
            raise

    def _chunk_text(self, text: str, doc_id: str, chapters: List[str], chunk_size: int = 800, overlap: int = 150) -> List[DocumentChunk]:
        """
        Splits text into overlapping semantic chunks with chapter association.
        """
        paragraphs = text.split("\n\n")
        chunks: List[DocumentChunk] = []
        current_chunk = ""
        current_chapter = chapters[0] if chapters else "General"
        page_num = 1
        chunk_idx = 0

        for para in paragraphs:
            clean_p = para.strip()
            if not clean_p:
                continue

            # Detect page numbers
            page_match = re.search(r"--- Page (\d+) ---", clean_p)
            if page_match:
                page_num = int(page_match.group(1))

            # Detect chapter changes
            for ch in chapters:
                if ch.lower() in clean_p.lower():
                    current_chapter = ch
                    break

            if len(current_chunk) + len(clean_p) > chunk_size:
                if current_chunk:
                    chunks.append(DocumentChunk(
                        chunk_id=f"{doc_id}_c{chunk_idx}",
                        doc_id=doc_id,
                        chapter=current_chapter,
                        content=current_chunk,
                        page_number=page_num
                    ))
                    chunk_idx += 1
                    # Overlap
                    current_chunk = current_chunk[-overlap:] + "\n" + clean_p
                else:
                    current_chunk = clean_p
            else:
                current_chunk += "\n" + clean_p if current_chunk else clean_p

        if current_chunk:
            chunks.append(DocumentChunk(
                chunk_id=f"{doc_id}_c{chunk_idx}",
                doc_id=doc_id,
                chapter=current_chapter,
                content=current_chunk,
                page_number=page_num
            ))

        return chunks

    def retrieve_context(self, doc_id: str, query: str, top_k: int = 4) -> str:
        """
        Retrieves relevant contextual chunks for a given query from an uploaded document.
        Uses keyword relevance and chapter grounding.
        """
        chunks = self.in_memory_chunks.get(doc_id, [])
        if not chunks:
            return ""

        query_words = set(re.findall(r"\w+", query.lower()))
        scored_chunks = []

        for chunk in chunks:
            chunk_words = set(re.findall(r"\w+", chunk.content.lower()))
            overlap_score = len(query_words.intersection(chunk_words))
            # Boost score if query matches chapter name
            if any(qw in chunk.chapter.lower() for qw in query_words):
                overlap_score += 2
            scored_chunks.append((overlap_score, chunk))

        scored_chunks.sort(key=lambda x: x[0], reverse=True)
        top_chunks = [c.content for _, c in scored_chunks[:top_k]]
        return "\n\n---\n\n".join(top_chunks)

    def synthesize_topic_curriculum(self, topic: str, level: str, language: str) -> str:
        """
        Synthesizes a grounded knowledge base for a topic when no document is uploaded.
        Simulates standard academic textbook grounding for the requested level and language.
        """
        return f"""
KNOWLEDGE BASE FOR TOPIC: {topic}
TARGET LEARNER LEVEL: {level}
TEACHING LANGUAGE: {language}

1. CORE DEFINITION:
{topic} is a fundamental concept requiring clear conceptual intuition before formal mathematical/technical definitions.

2. INTUITIVE ANALOGIES (For {level}):
- If Beginner: Use everyday physical metaphors (e.g. water flow, moving cars, kitchen recipes).
- If Intermediate: Step-by-step causality, input-output transformations.
- If Advanced: Formal governing laws, mathematical bounds, optimization tradeoffs.

3. GOVERNING PRINCIPLES & FORMULAS:
- Key relationships and foundational laws governing {topic}.
- Step-by-step visual demonstration needed on the blackboard.

4. COMMON STUDENT MISCONCEPTIONS TO TARGET:
- Conflating rate of change with current state.
- Inverting dependent and independent variables.
- Misapplying linear models to non-linear systems.
"""

# Global singleton instance
rag_engine = RAGEngine()
